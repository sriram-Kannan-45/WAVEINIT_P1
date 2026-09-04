import tempfile
from urllib.parse import urlparse
from pathlib import Path
from typing import Optional

import requests
from bs4 import BeautifulSoup

from .config import RAGConfig


class UnsupportedSourceError(ValueError):
    pass


class TextExtractor:
    def __init__(self, config: RAGConfig):
        self.config = config

    def validate_file(self, file_path: str, mime_type: Optional[str] = None) -> Path:
        path = Path(file_path).resolve()
        if not path.exists() or not path.is_file():
            raise FileNotFoundError(f"File not found: {path}")
        if path.stat().st_size > self.config.max_file_size_bytes:
            raise UnsupportedSourceError("File too large for AI quiz generation.")
        suffix = path.suffix.lower()
        if suffix not in self.config.allowed_extensions:
            raise UnsupportedSourceError(f"Unsupported file type: {suffix}")
        if mime_type and mime_type not in self.config.allowed_mimes:
            if suffix not in self.config.allowed_extensions:
                raise UnsupportedSourceError(f"Unsupported MIME type: {mime_type}")
        return path

    def extract_from_file(self, file_path: str, mime_type: Optional[str] = None) -> str:
        path = self.validate_file(file_path, mime_type)
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self._extract_pdf(path)
        if suffix == ".docx":
            return self._extract_docx(path)
        if suffix == ".pptx":
            return self._extract_pptx(path)
        if suffix == ".txt":
            return self._extract_txt(path)
        raise UnsupportedSourceError(f"Unsupported file type: {suffix}")

    def extract_from_url(self, url: str) -> str:
        if not url or not url.lower().startswith(("http://", "https://")):
            raise UnsupportedSourceError("URL must start with http:// or https://")
        with requests.get(url, timeout=20, stream=True, headers={"User-Agent": "LMS-AI-Quiz-RAG/1.0"}) as response:
            response.raise_for_status()
            content_type = response.headers.get("content-type", "").split(";")[0].strip().lower()
            data = bytearray()
            for chunk in response.iter_content(chunk_size=65536):
                data.extend(chunk)
                if len(data) > self.config.max_file_size_bytes:
                    raise UnsupportedSourceError("Source URL exceeds the maximum file size.")
            mime_extensions = {"application/pdf": ".pdf", "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx", "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx", "text/plain": ".txt"}
            suffix = mime_extensions.get(content_type) or Path(urlparse(response.url).path).suffix.lower()
            if suffix in self.config.allowed_extensions and content_type not in {"text/html", "application/xhtml+xml"}:
                with tempfile.TemporaryDirectory(prefix="quiz-source-") as temp_dir:
                    source = Path(temp_dir) / ("source" + suffix)
                    source.write_bytes(data)
                    return self.extract_from_file(str(source))
            if content_type not in {"text/html", "application/xhtml+xml"}:
                raise UnsupportedSourceError(f"Unsupported URL content type: {content_type}")
            html = bytes(data).decode(response.encoding or "utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript", "svg", "header", "footer", "nav", "aside"]):
            tag.decompose()
        main = soup.find("main") or soup.find("article") or soup.body or soup
        return main.get_text("\n", strip=True)

    @staticmethod
    def _extract_pdf(path: Path) -> str:
        import fitz

        parts = []
        with fitz.open(path) as document:
            for page in document:
                parts.append(page.get_text("text"))
        return "\n\n".join(parts)

    @staticmethod
    def _extract_docx(path: Path) -> str:
        import docx

        document = docx.Document(path)
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)

    @staticmethod
    def _extract_pptx(path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(path)
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines = [f"Slide {index}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
            if len(lines) > 1:
                slides.append("\n".join(lines))
        return "\n\n".join(slides)

    @staticmethod
    def _extract_txt(path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

